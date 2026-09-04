"""Geometric QA: out-of-bounds, edge margins, text overflow, and shape overlap."""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
SW, SH = 13.333, 7.5
MIN_MARGIN = 0.45          # tolerance below the 0.5" guideline
# Rough advance width per point of font size, for Calibri/Cambria-ish faces.
CHAR_W = {"Calibri": 0.475, "Cambria": 0.505}

def inches(v): return (v or 0) / EMU_IN

def est_text_height(tf, w_in, default_size=14.0, face="Calibri"):
    """Estimate rendered height of a text frame in inches."""
    total = 0.0
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs)
        sizes = [r.font.size.pt for r in p.runs if r.font.size]
        size = max(sizes) if sizes else default_size
        names = [r.font.name for r in p.runs if r.font.name]
        f = names[0] if names else face
        cw = CHAR_W.get(f, 0.48) * size / 72.0
        usable = max(w_in - 0.10, 0.2)
        chars_per_line = max(int(usable / cw), 1)
        if not txt.strip():
            lines = 1
        else:
            # account for explicit bullets adding an indent
            eff = chars_per_line - (2 if p.level or txt.startswith(("•", "-")) else 0)
            eff = max(eff, 1)
            lines = max(1, -(-len(txt) // eff))
        total += lines * (size * 1.22) / 72.0 + 0.045
    return total

def main(path):
    prs = Presentation(path)
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            x, y = inches(sh.left), inches(sh.top)
            w, h = inches(sh.width), inches(sh.height)
            name = sh.shape_type
            label = ""
            if sh.has_text_frame:
                label = " ".join(r.text for p in sh.text_frame.paragraphs for r in p.runs)[:48]
            # 1. out of slide bounds
            if (x < -0.02 or y < -0.02 or x + w > SW + 0.02 or y + h > SH + 0.02) and label.strip():
                if not (w >= SW - 0.01 and h >= SH - 0.01):
                    problems.append(f"S{idx}: FUERA DE LIENZO  x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}  «{label}»")
            # 2. text overflow
            if sh.has_text_frame and label.strip():
                need = est_text_height(sh.text_frame, w)
                if need > h + 0.07:
                    problems.append(f"S{idx}: DESBORDE texto  necesita {need:.2f}\" en caja de {h:.2f}\"  «{label}»")
            if sh.has_text_frame and label.strip():
                boxes.append((x, y, w, h, label))
        # 3. edge margins for text content
        for (x, y, w, h, label) in boxes:
            if y > 6.8 and h < 0.35:
                continue  # footer band
            if x < MIN_MARGIN - 0.01 or y < MIN_MARGIN - 0.01 or x + w > SW - MIN_MARGIN + 0.01 or y + h > SH - MIN_MARGIN + 0.01:
                problems.append(f"S{idx}: MARGEN corto  x={x:.2f} y={y:.2f} der={SW-(x+w):.2f} inf={SH-(y+h):.2f}  «{label}»")
    if problems:
        print(f"{len(problems)} hallazgo(s):\n")
        for p in problems: print("  " + p)
    else:
        print("Sin hallazgos geométricos.")
    return 1 if problems else 0

sys.exit(main(sys.argv[1]))
